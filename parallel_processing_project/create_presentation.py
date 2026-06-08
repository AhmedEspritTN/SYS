"""Generate a styled PowerPoint with aligned diagrams, images, and navigation."""

from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# Layout grid (10" x 7.5" slide)
SLIDE_W = 10.0
SLIDE_H = 7.5
MARGIN_X = 0.55
CONTENT_W = SLIDE_W - 2 * MARGIN_X
HEADER_H = 1.0
FOOTER_TOP = 7.05
BODY_TOP = 1.35

THEME = {
    "navy": "#1E3A5F",
    "teal": "#0D9488",
    "sky": "#38BDF8",
    "green": "#059669",
    "bg": "#F8FAFC",
    "card": "#FFFFFF",
    "muted": "#64748B",
    "text": "#1E293B",
    "light_teal": "#CCFBF1",
    "light_sky": "#E0F2FE",
}

COLORS = {
    "navy": RGBColor(0x1E, 0x3A, 0x5F),
    "teal": RGBColor(0x0D, 0x94, 0x88),
    "sky": RGBColor(0x38, 0xBD, 0xF8),
    "green": RGBColor(0x05, 0x96, 0x69),
    "bg": RGBColor(0xF8, 0xFA, 0xFC),
    "card": RGBColor(0xFF, 0xFF, 0xFF),
    "muted": RGBColor(0x64, 0x74, 0x8B),
    "text": RGBColor(0x1E, 0x29, 0x3B),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
}

AUTHORS = ["Hassen Moussi", "Ahmed Arfaoui"]


def hex_rgb(name: str) -> RGBColor:
    value = THEME[name].lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


class Deck:
    """Build slides and keep references for hyperlinks."""

    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.slides: dict[str, object] = {}
        self.assets = Path(__file__).parent / "presentation_assets"
        self.assets.mkdir(exist_ok=True)

    def blank(self, key: str):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slides[key] = slide
        return slide


def set_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_round(slide, left, top, width, height, fill, line=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape


def add_header(slide, title: str, subtitle: str = "") -> None:
    add_rect(slide, 0, 0, SLIDE_W, HEADER_H, COLORS["navy"])
    add_rect(slide, 0, HEADER_H, SLIDE_W, 0.07, COLORS["teal"]).line.fill.background()

    box = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(0.15), Inches(CONTENT_W), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(0.58), Inches(CONTENT_W), Inches(0.35))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(12)
        sp.font.color.rgb = COLORS["sky"]


def add_footer(slide, text: str = "IntegrityCheck  |  Hassen Moussi & Ahmed Arfaoui") -> None:
    add_rect(slide, 0, FOOTER_TOP, SLIDE_W, 0.45, COLORS["navy"]).line.fill.background()
    box = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(FOOTER_TOP + 0.08), Inches(CONTENT_W), Inches(0.3))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = COLORS["sky"]
    p.alignment = PP_ALIGN.CENTER


def add_picture_fit(slide, image_path: Path, left: float, top: float, width: float, height: float) -> None:
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), Inches(width), Inches(height))


def add_nav_button(slide, deck: Deck, label: str, left: float, top: float, width: float, target_key: str) -> None:
    btn = add_round(slide, left, top, width, 0.42, COLORS["teal"])
    frame = btn.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]
    p.alignment = PP_ALIGN.CENTER
    if target_key in deck.slides:
        btn.click_action.target_slide = deck.slides[target_key]


def add_bullets(slide, bullets: list[str], top: float = BODY_TOP, width: float = CONTENT_W, size: int = 16):
    box = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(top), Inches(width), Inches(FOOTER_TOP - top - 0.2))
    frame = box.text_frame
    frame.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        p.text = bullet
        p.font.size = Pt(size - 2 if bullet.startswith("  ") else size)
        p.font.color.rgb = COLORS["muted"] if bullet.startswith("  ") else COLORS["text"]
        p.space_after = Pt(8)
        p.level = 1 if bullet.startswith("  ") else 0


# ── Matplotlib diagram generators ──────────────────────────────────────────

def _save_fig(path: Path) -> None:
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight", facecolor=THEME["bg"])
    plt.close()


def generate_pipeline_image(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(11, 2.6))
    ax.set_xlim(0, 11)
    ax.set_ylim(0, 3)
    ax.axis("off")
    ax.set_facecolor(THEME["bg"])

    steps = [
        ("User", "Select file", THEME["light_sky"], THEME["navy"]),
        ("Client", "client.py", THEME["card"], THEME["navy"]),
        ("Split", "Chunks", THEME["light_teal"], THEME["teal"]),
        ("Queue", "Waiting room", THEME["card"], THEME["teal"]),
        ("Workers", "Threads / Processes", THEME["light_sky"], THEME["navy"]),
        ("Hash", "SHA-256", THEME["light_teal"], THEME["teal"]),
        ("Verify", "Report ✓", THEME["green"], "white"),
    ]

    x = 0.2
    box_w = 1.35
    for i, (title, sub, bg, fg) in enumerate(steps):
        box = FancyBboxPatch(
            (x, 0.9), box_w, 1.3, boxstyle="round,pad=0.08,rounding_size=0.15",
            facecolor=bg, edgecolor=THEME["muted"], linewidth=1.2,
        )
        ax.add_patch(box)
        ax.text(x + box_w / 2, 1.75, title, ha="center", va="center", fontsize=10, fontweight="bold", color=fg)
        ax.text(x + box_w / 2, 1.25, sub, ha="center", va="center", fontsize=9, color=fg)
        if i < len(steps) - 1:
            ax.annotate("", xy=(x + box_w + 0.18, 1.55), xytext=(x + box_w + 0.02, 1.55),
                        arrowprops=dict(arrowstyle="-|>", color=THEME["teal"], lw=2))
        x += box_w + 0.28

    ax.text(5.5, 0.25, "IntegrityCheck Processing Pipeline", ha="center", fontsize=12,
            fontweight="bold", color=THEME["navy"])
    _save_fig(path)


def generate_barber_image(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(10, 3.2))
    ax.set_xlim(0, 10)
    ax.set_ylim(0, 3.5)
    ax.axis("off")
    ax.set_facecolor(THEME["bg"])

    boxes = [
        (0.3, "File Chunks\n(customers)", THEME["light_sky"], 1.5),
        (2.4, "Waiting Room\n(8 chairs max)", THEME["light_teal"], 2.0),
        (5.0, "Barber Threads\n(workers)", THEME["card"], 2.0),
        (7.6, "SHA-256\nVerified ✓", THEME["green"], 1.6),
    ]
    for x, label, color, w in boxes:
        box = FancyBboxPatch(
            (x, 1.0), w, 1.5, boxstyle="round,pad=0.06,rounding_size=0.12",
            facecolor=color, edgecolor=THEME["navy"], linewidth=1.2,
        )
        ax.add_patch(box)
        fg = "white" if color == THEME["green"] else THEME["navy"]
        ax.text(x + w / 2, 1.75, label, ha="center", va="center", fontsize=10, fontweight="bold", color=fg)

    for x1, x2 in [(1.8, 2.4), (4.4, 5.0), (7.0, 7.6)]:
        ax.add_patch(FancyArrowPatch((x1, 1.75), (x2, 1.75), arrowstyle="-|>",
                                     color=THEME["teal"], lw=2.5, mutation_scale=14))

    ax.text(5.0, 0.35, "Sleeping Barber — chunks wait in chairs, barbers process in parallel",
            ha="center", fontsize=11, fontweight="bold", color=THEME["navy"])
    _save_fig(path)


def generate_architecture_image(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(8, 3.5))
    ax.set_xlim(0, 8)
    ax.set_ylim(0, 4)
    ax.axis("off")
    ax.set_facecolor(THEME["bg"])

    layers = [
        ("User Layer", "client.py", THEME["sky"], 3.0),
        ("Processing Layer", "file_processing  •  sleeping_barber  •  parallel_processing", THEME["teal"], 2.4),
        ("OS Concepts", "ipc_communication  •  synchronization", THEME["navy"], 2.4),
    ]
    y = 3.2
    for title, modules, color, h in layers:
        box = FancyBboxPatch(
            (1.0, y - h), 6.0, h, boxstyle="round,pad=0.05,rounding_size=0.1",
            facecolor=color, edgecolor=THEME["navy"], linewidth=1.2,
        )
        ax.add_patch(box)
        fg = THEME["navy"] if color == THEME["sky"] else "white"
        ax.text(4.0, y - h / 2 + 0.2, title, ha="center", fontsize=11, fontweight="bold", color=fg)
        ax.text(4.0, y - h / 2 - 0.25, modules, ha="center", fontsize=8.5, color=fg)
        y -= h + 0.25
        if y > 0.5:
            ax.annotate("", xy=(4.0, y + 0.12), xytext=(4.0, y + 0.22),
                        arrowprops=dict(arrowstyle="-|>", color=THEME["muted"], lw=1.8))

    _save_fig(path)


def generate_client_mock_image(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(5.5, 4.5))
    ax.set_xlim(0, 5.5)
    ax.set_ylim(0, 4.5)
    ax.axis("off")
    ax.set_facecolor(THEME["navy"])

    # Terminal window
    win = FancyBboxPatch(
        (0.3, 0.3), 4.9, 3.9, boxstyle="round,pad=0.04,rounding_size=0.08",
        facecolor="#0F172A", edgecolor=THEME["teal"], linewidth=2,
    )
    ax.add_patch(win)
    ax.text(0.55, 3.85, "● ● ●", fontsize=10, color=THEME["teal"], family="monospace")
    ax.text(0.55, 3.45, "PARALLEL FILE PROCESSING - CLIENT", fontsize=9, color=THEME["sky"], fontweight="bold")

    menu = [
        "1. Process file (Sleeping Barber)",
        "2. Process file (Multiprocessing)",
        "3. Process file (Multithreading)",
        "4. Benchmark all modes",
        "5. Load external solution",
        "6. Exit",
        "",
        "Your choice: 1",
        "File: sample_input.bin",
        "Barbers: 4  |  Chairs: 8",
        "Processing chunk 0... done ✓",
    ]
    y = 3.05
    for line in menu:
        color = THEME["green"] if "done" in line else ("white" if line.startswith("Your") or line.startswith("File") or "Barbers" in line else THEME["light_sky"])
        weight = "bold" if line.startswith("Your") else "normal"
        ax.text(0.55, y, line, fontsize=8.5, color=color, family="monospace", fontweight=weight)
        y -= 0.24

    ax.text(2.75, 0.15, "Interactive client.py", ha="center", fontsize=9, color=THEME["sky"])
    _save_fig(path)


def generate_benchmark_image(path: Path) -> None:
    fig, ax = plt.subplots(figsize=(6, 3.5))
    modes = ["Sequential", "Multithreading", "Sleeping Barber", "Multiprocessing"]
    times = [0.05, 0.03, 0.11, 0.21]
    colors = [THEME["muted"], THEME["teal"], THEME["navy"], THEME["sky"]]

    bars = ax.bar(modes, times, color=colors, edgecolor=THEME["navy"], linewidth=0.8)
    ax.set_ylabel("Time (seconds)", fontsize=10)
    ax.set_title("Benchmark — 20 MB file (lower is better)", fontsize=11, fontweight="bold", color=THEME["navy"])
    ax.set_facecolor(THEME["bg"])
    fig.patch.set_facecolor(THEME["bg"])
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)

    for bar, t in zip(bars, times):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.008,
                f"{t:.2f}s", ha="center", fontsize=9, fontweight="bold")

    plt.xticks(rotation=12, fontsize=9)
    _save_fig(path)


def generate_all_assets(deck: Deck) -> dict[str, Path]:
    paths = {
        "pipeline": deck.assets / "pipeline.png",
        "barber": deck.assets / "sleeping_barber.png",
        "architecture": deck.assets / "architecture.png",
        "client": deck.assets / "client_mock.png",
        "benchmark": deck.assets / "benchmark.png",
    }
    generate_pipeline_image(paths["pipeline"])
    generate_barber_image(paths["barber"])
    generate_architecture_image(paths["architecture"])
    generate_client_mock_image(paths["client"])
    generate_benchmark_image(paths["benchmark"])
    return paths


# ── Slides ─────────────────────────────────────────────────────────────────

def add_title_slide(deck: Deck) -> None:
    slide = deck.blank("title")
    set_background(slide, COLORS["navy"])
    add_rect(slide, 0, 0, 0.3, SLIDE_H, COLORS["teal"]).line.fill.background()
    add_rect(slide, 0, 6.85, SLIDE_W, 0.08, COLORS["teal"]).line.fill.background()

    # Logo
    shield = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(0.55), Inches(0.45), Inches(0.9), Inches(0.9))
    shield.fill.solid()
    shield.fill.fore_color.rgb = COLORS["teal"]
    shield.line.fill.background()

    logo_text = slide.shapes.add_textbox(Inches(0.72), Inches(0.68), Inches(0.6), Inches(0.4))
    logo_text.text_frame.paragraphs[0].text = "✓"
    logo_text.text_frame.paragraphs[0].font.size = Pt(22)
    logo_text.text_frame.paragraphs[0].font.bold = True
    logo_text.text_frame.paragraphs[0].font.color.rgb = COLORS["white"]

    # Title block — left aligned in content area
    title = slide.shapes.add_textbox(Inches(0.55), Inches(1.6), Inches(5.5), Inches(0.9))
    p = title.text_frame.paragraphs[0]
    p.text = "IntegrityCheck"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS["white"]

    sub = slide.shapes.add_textbox(Inches(0.55), Inches(2.45), Inches(5.8), Inches(0.6))
    sub.text_frame.paragraphs[0].text = "Parallel File Integrity Verification System"
    sub.text_frame.paragraphs[0].font.size = Pt(18)
    sub.text_frame.paragraphs[0].font.color.rgb = COLORS["sky"]

    course = slide.shapes.add_textbox(Inches(0.55), Inches(3.15), Inches(5.5), Inches(0.9))
    course.text_frame.paragraphs[0].text = (
        "Systèmes d'Exploitation Avancé\nParallel Processing & Benchmarking"
    )
    course.text_frame.paragraphs[0].font.size = Pt(14)
    course.text_frame.paragraphs[0].font.color.rgb = COLORS["white"]

    # Authors
    authors_box = add_round(slide, 0.55, 4.15, 4.8, 0.95, COLORS["teal"])
    af = authors_box.text_frame
    af.clear()
    af.vertical_anchor = MSO_ANCHOR.MIDDLE
    ap1 = af.paragraphs[0]
    ap1.text = "Presented by"
    ap1.font.size = Pt(10)
    ap1.font.color.rgb = COLORS["white"]
    ap1.alignment = PP_ALIGN.CENTER
    ap2 = af.add_paragraph()
    ap2.text = f"{AUTHORS[0]}  &  {AUTHORS[1]}"
    ap2.font.size = Pt(16)
    ap2.font.bold = True
    ap2.font.color.rgb = COLORS["white"]
    ap2.alignment = PP_ALIGN.CENTER

    # Right side — client mock image
    client_img = deck.assets / "client_mock.png"
    if client_img.exists():
        add_picture_fit(slide, client_img, 5.55, 0.55, 3.9, 3.55)

    # Interactive navigation buttons
    nav_y = 5.35
    btn_w = 2.05
    gap = 0.15
    labels = [
        ("→ Pipeline", "pipeline"),
        ("→ Sleeping Barber", "barber"),
        ("→ Architecture", "architecture"),
        ("→ Demo", "demo"),
    ]
    x = 0.55
    for label, key in labels:
        add_nav_button(slide, deck, label, x, nav_y, btn_w, key)
        x += btn_w + gap

    tag = add_round(slide, 0.55, 6.05, 4.8, 0.5, hex_rgb("light_teal"), COLORS["teal"])
    tag.text_frame.paragraphs[0].text = "Click buttons to jump  •  SHA-256  •  Multiprocessing  •  Sleeping Barber"
    tag.text_frame.paragraphs[0].font.size = Pt(9)
    tag.text_frame.paragraphs[0].font.color.rgb = COLORS["navy"]
    tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_image_slide(deck: Deck, key: str, title: str, subtitle: str, image_key: str, assets: dict, caption: str = ""):
    slide = deck.blank(key)
    set_background(slide, COLORS["bg"])
    add_header(slide, title, subtitle)
    add_footer(slide)

    img_path = assets[image_key]
    img_top = BODY_TOP + 0.1
    img_h = 3.35 if caption else 4.0
    add_picture_fit(slide, img_path, MARGIN_X, img_top, CONTENT_W, img_h)

    if caption:
        cap_box = add_round(slide, MARGIN_X, img_top + img_h + 0.2, CONTENT_W, FOOTER_TOP - img_top - img_h - 0.35, COLORS["card"], COLORS["teal"])
        cap = cap_box.text_frame
        cap.clear()
        cap.word_wrap = True
        cp = cap.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(12)
        cp.font.color.rgb = COLORS["text"]
        cp.alignment = PP_ALIGN.LEFT


def add_modes_slide(deck: Deck) -> None:
    slide = deck.blank("modes")
    set_background(slide, COLORS["bg"])
    add_header(slide, "Processing Modes", "Four parallel strategies — same SHA-256 result")
    add_footer(slide)

    modes = [
        ("Sequential", "1 worker", THEME["muted"], COLORS["white"]),
        ("Multiprocessing", "N processes", THEME["sky"], COLORS["navy"]),
        ("Multithreading", "N threads", THEME["teal"], COLORS["white"]),
        ("Sleeping Barber", "Bounded queue", THEME["navy"], COLORS["white"]),
    ]
    card_w = 2.0
    total_w = len(modes) * card_w
    gap = (CONTENT_W - total_w) / (len(modes) - 1)
    x = MARGIN_X
    y = BODY_TOP + 0.15
    card_h = 1.35

    for name, desc, hex_color, text_color in modes:
        fill = hex_rgb("muted") if hex_color == THEME["muted"] else RGBColor(
            int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16)
        )
        card = add_round(slide, x, y, card_w, card_h, fill)
        frame = card.text_frame
        frame.clear()
        frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = frame.paragraphs[0]
        p1.text = name
        p1.font.bold = True
        p1.font.size = Pt(13)
        p1.font.color.rgb = text_color
        p1.alignment = PP_ALIGN.CENTER
        p2 = frame.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = text_color
        p2.alignment = PP_ALIGN.CENTER
        x += card_w + gap

    # Benchmark image centered below
    bench = deck.assets / "benchmark.png"
    if bench.exists():
        add_picture_fit(slide, bench, MARGIN_X + 0.8, y + card_h + 0.35, CONTENT_W - 1.6, 2.85)


def add_content_slide(deck: Deck, key: str, title: str, bullets: list[str], subtitle: str = "") -> None:
    slide = deck.blank(key)
    set_background(slide, COLORS["bg"])
    add_header(slide, title, subtitle)
    add_bullets(slide, bullets)
    add_footer(slide)


def add_demo_slide(deck: Deck, assets: dict) -> None:
    slide = deck.blank("demo")
    set_background(slide, COLORS["bg"])
    add_header(slide, "Live Demo", "How to run IntegrityCheck")
    add_footer(slide)

    left_w = 4.1
    add_picture_fit(slide, assets["client"], MARGIN_X, BODY_TOP, left_w, 3.5)

    cmds = add_round(slide, MARGIN_X + left_w + 0.25, BODY_TOP, CONTENT_W - left_w - 0.25, 3.5, COLORS["card"], COLORS["teal"])
    lines = [
        "Commands to demo:",
        "",
        "python client.py",
        "  → Interactive menu (recommended)",
        "",
        "python parallel_processing.py sample_input.bin",
        "  --mode sleeping_barber",
        "",
        "python parallel_processing.py sample_input.bin",
        "  --mode benchmark",
        "  --output-report processing_report.json",
    ]
    tf = cmds.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13 if i == 0 else 11)
        p.font.bold = i == 0
        p.font.name = "Consolas" if line.startswith("python") else None
        p.font.color.rgb = COLORS["navy"] if i == 0 else COLORS["text"]
        p.space_after = Pt(3)

    note = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(5.0), Inches(CONTENT_W), Inches(1.5))
    np = note.text_frame.paragraphs[0]
    np.text = "Tip for demo video: show option 1 (Sleeping Barber) with shop-full messages, then option 4 (benchmark)."
    np.font.size = Pt(12)
    np.font.color.rgb = COLORS["muted"]
    np.font.italic = True


def add_closing_slide(deck: Deck) -> None:
    slide = deck.blank("thanks")
    set_background(slide, COLORS["navy"])
    add_rect(slide, 0, 3.15, SLIDE_W, 0.06, COLORS["teal"]).line.fill.background()

    t = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(2.0), Inches(CONTENT_W), Inches(0.9))
    t.text_frame.paragraphs[0].text = "Thank You"
    t.text_frame.paragraphs[0].font.size = Pt(38)
    t.text_frame.paragraphs[0].font.bold = True
    t.text_frame.paragraphs[0].font.color.rgb = COLORS["white"]
    t.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    s = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(3.45), Inches(CONTENT_W), Inches(1.2))
    s.text_frame.paragraphs[0].text = "IntegrityCheck — Fast • Parallel • Verifiable"
    s.text_frame.paragraphs[0].font.size = Pt(18)
    s.text_frame.paragraphs[0].font.color.rgb = COLORS["sky"]
    s.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    a = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(4.5), Inches(CONTENT_W), Inches(0.8))
    a.text_frame.paragraphs[0].text = f"{AUTHORS[0]}  &  {AUTHORS[1]}"
    a.text_frame.paragraphs[0].font.size = Pt(16)
    a.text_frame.paragraphs[0].font.bold = True
    a.text_frame.paragraphs[0].font.color.rgb = COLORS["white"]
    a.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    q = slide.shapes.add_textbox(Inches(MARGIN_X), Inches(5.4), Inches(CONTENT_W), Inches(0.5))
    q.text_frame.paragraphs[0].text = "Questions?"
    q.text_frame.paragraphs[0].font.size = Pt(20)
    q.text_frame.paragraphs[0].font.color.rgb = COLORS["sky"]
    q.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


def build_presentation() -> Presentation:
    deck = Deck()
    assets = generate_all_assets(deck)

    # Build title first (needs assets + slide refs for later buttons — buttons added at end)
    # Generate client image before title
    add_title_slide(deck)  # nav buttons reference slides not yet created — fix order below

    add_content_slide(deck, "context", "Project Context", [
        "Course: Systèmes d'Exploitation Avancé",
        "Topic: high-performance parallel processing & benchmarking",
        "Application: large file processing (team choice)",
        "Product: IntegrityCheck — verify files after transfer or backup",
        "Tech: Python 3.8+, standard library only",
    ])

    add_content_slide(deck, "problem", "The Problem", [
        "IT teams transfer large backup files and archives daily.",
        "A corrupted file can break deployments and backups.",
        "Sequential checksum tools waste multi-core CPU power.",
        "We need fast, parallel, auditable file verification.",
    ])

    add_image_slide(
        deck, "pipeline", "Processing Pipeline", "How a file moves through IntegrityCheck",
        "pipeline", assets,
        "User selects a file → client splits it into chunks → chunks enter the waiting queue → "
        "workers hash in parallel → JSON report confirms integrity.",
    )

    add_image_slide(
        deck, "barber", "Sleeping Barber", "Classic OS algorithm applied to file chunks",
        "barber", assets,
        "Customer = chunk  |  Barber = thread  |  Chair = queue slot  |  Semaphore wakes barbers.",
    )

    add_modes_slide(deck)

    add_image_slide(
        deck, "architecture", "Project Architecture", "Three clear layers",
        "architecture", assets,
        "Client layer for users, processing layer for parallel file work, OS concepts layer for IPC and sync.",
    )

    add_content_slide(deck, "ipc", "IPC & Synchronization", [
        "Inter-process communication modules:",
        "  Pipes — one-to-one process communication",
        "  Queues — multi-producer / multi-consumer",
        "  Shared memory — protected shared data",
        "Synchronization: semaphores, Dining Philosophers, Sleeping Barber, Producer-Consumer",
    ])

    add_content_slide(deck, "software", "Software Solution Loading", [
        "External Python modules loaded at runtime via importlib.",
        "Plugin example: sample_solution.py",
        "Configurable threads and chunk size.",
        "Enables porting existing file-processing code.",
    ])

    add_demo_slide(deck, assets)

    add_content_slide(deck, "results", "Results & Deliverables", [
        "Metrics: time, throughput (MB/s), speedup, SHA-256",
        "Output: processing_report.json",
        "Deliverables: source code, REQUIREMENTS.md, demo video, presentation",
        f"Authors: {AUTHORS[0]}, {AUTHORS[1]}",
    ])

    add_closing_slide(deck)

    # Re-wire title nav buttons now that all slides exist
    title_slide = deck.slides["title"]
    for shape in title_slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.startswith("→"):
            label = shape.text_frame.text
            mapping = {
                "→ Pipeline": "pipeline",
                "→ Sleeping Barber": "barber",
                "→ Architecture": "architecture",
                "→ Demo": "demo",
            }
            if label in mapping and mapping[label] in deck.slides:
                shape.click_action.target_slide = deck.slides[mapping[label]]

    return deck.prs


def main() -> None:
    folder = Path(__file__).parent
    output = folder / "presentation.pptx"
    prs = build_presentation()

    try:
        prs.save(output)
    except PermissionError:
        output = folder / "presentation_styled.pptx"
        prs.save(output)
        print("Note: close presentation.pptx in PowerPoint, then re-run to overwrite.")

    print(f"Presentation saved to: {output}")
    print(f"Diagram images saved to: {folder / 'presentation_assets'}")


if __name__ == "__main__":
    main()
